"""
Document loader — converts files into plain text.

Supported: .txt, .md, .pdf, .docx, .xlsx, .xls, .csv, .pptx
Each loader returns a list of page/sheet dicts: {text, source, page}
"""

import os
from pathlib import Path
from typing import List, Dict, Any

from .exceptions import LoaderError


def load_file(file_path: str) -> List[Dict[str, Any]]:
    """Load a file and return a list of {text, source, page} dicts."""
    path = Path(file_path)

    if not path.exists():
        raise LoaderError(f"File not found: {file_path}")
    if not os.access(file_path, os.R_OK):
        raise LoaderError(f"Permission denied: {file_path}")
    if path.stat().st_size == 0:
        raise LoaderError(f"File is empty: {path.name}")

    ext = path.suffix.lower()
    loaders = {
        ".txt":  _load_text,
        ".md":   _load_text,
        ".pdf":  _load_pdf,
        ".docx": _load_docx,
        ".doc":  _load_docx,
        ".xlsx": _load_spreadsheet,
        ".xls":  _load_spreadsheet,
        ".csv":  _load_csv,
        ".pptx": _load_pptx,
    }

    loader = loaders.get(ext)
    if loader is None:
        raise LoaderError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(sorted(loaders))}"
        )

    try:
        pages = loader(str(path))
    except LoaderError:
        raise
    except Exception as exc:
        raise LoaderError(
            f"Failed to parse '{path.name}': {exc}"
        ) from exc

    # Attach source filename and drop pages with no usable text
    result = []
    for p in pages:
        p["source"] = path.name
        if p.get("text", "").strip():
            result.append(p)

    if not result:
        raise LoaderError(
            f"No readable text could be extracted from '{path.name}'. "
            "The file may be scanned/image-only, encrypted, or empty."
        )

    return result


# ── individual loaders ────────────────────────────────────────────────────────

def _load_text(path: str) -> List[Dict[str, Any]]:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return [{"text": f.read(), "page": 1}]
    except OSError as exc:
        raise LoaderError(f"Cannot read text file: {exc}") from exc


def _load_pdf(path: str) -> List[Dict[str, Any]]:
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError, PdfStreamError
    except ImportError as exc:
        raise LoaderError("pypdf is not installed. Run: pip install pypdf") from exc

    try:
        reader = PdfReader(path)
    except Exception as exc:
        raise LoaderError(f"Cannot open PDF (possibly encrypted or corrupted): {exc}") from exc

    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""   # skip unreadable pages silently
        if text.strip():
            pages.append({"text": text, "page": i})
    return pages


def _load_docx(path: str) -> List[Dict[str, Any]]:
    try:
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as exc:
        raise LoaderError("python-docx is not installed. Run: pip install python-docx") from exc

    try:
        doc = Document(path)
    except Exception as exc:
        raise LoaderError(f"Cannot open DOCX (possibly corrupted or wrong format): {exc}") from exc

    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": text, "page": 1}]


def _load_spreadsheet(path: str) -> List[Dict[str, Any]]:
    try:
        import pandas as pd
    except ImportError as exc:
        raise LoaderError("pandas is not installed. Run: pip install pandas openpyxl") from exc

    try:
        xl = pd.ExcelFile(path)
    except Exception as exc:
        raise LoaderError(f"Cannot open spreadsheet (possibly corrupted or password-protected): {exc}") from exc

    pages = []
    for sheet in xl.sheet_names:
        try:
            df = xl.parse(sheet)
            if df.empty:
                continue
            text = f"Sheet: {sheet}\n{df.to_string(index=False)}"
            pages.append({"text": text, "page": sheet})
        except Exception:
            continue   # skip unreadable sheets, continue with others
    return pages


def _load_csv(path: str) -> List[Dict[str, Any]]:
    try:
        import pandas as pd
    except ImportError as exc:
        raise LoaderError("pandas is not installed. Run: pip install pandas") from exc

    try:
        df = pd.read_csv(path, encoding_errors="replace")
    except Exception as exc:
        raise LoaderError(f"Cannot parse CSV: {exc}") from exc

    if df.empty:
        raise LoaderError("CSV file contains no data rows.")

    return [{"text": df.to_string(index=False), "page": 1}]


def _load_pptx(path: str) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise LoaderError("python-pptx is not installed. Run: pip install python-pptx") from exc

    try:
        prs = Presentation(path)
    except Exception as exc:
        raise LoaderError(f"Cannot open PPTX (possibly corrupted): {exc}") from exc

    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            except Exception:
                continue
        if texts:
            pages.append({"text": "\n".join(texts), "page": i})
    return pages
